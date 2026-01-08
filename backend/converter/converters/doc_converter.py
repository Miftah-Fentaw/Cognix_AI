import io
import os
import tempfile
import subprocess
from docx import Document
from pptx import Presentation
from pptx.util import Pt

def convert(file, target_type):
    """
    Convert DOC/DOCX files to various formats.
    Supports: DOC/DOCX → PDF, TXT, PPTX
    """
    try:
        output = io.BytesIO()
        doc = Document(file)

        if target_type.upper() == 'PDF':
            # Use LibreOffice for DOCX to PDF conversion
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_docx:
                doc.save(temp_docx.name)
                temp_docx_path = temp_docx.name
            
            try:
                # Create temporary output directory
                temp_dir = tempfile.mkdtemp()
                
                # Use LibreOffice to convert to PDF
                subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', temp_dir, temp_docx_path
                ], check=True, capture_output=True, timeout=30)
                
                # Read the generated PDF
                pdf_filename = os.path.splitext(os.path.basename(temp_docx_path))[0] + '.pdf'
                pdf_path = os.path.join(temp_dir, pdf_filename)
                
                with open(pdf_path, 'rb') as pdf_file:
                    output.write(pdf_file.read())
                
                # Cleanup
                os.remove(pdf_path)
                os.rmdir(temp_dir)
            finally:
                os.remove(temp_docx_path)
        
        elif target_type.upper() == 'TXT':
            # Extract text from all paragraphs
            text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            output.write(text.encode('utf-8'))
        
        elif target_type.upper() == 'PPTX':
            # Convert DOCX to PPTX (each paragraph becomes a slide)
            prs = Presentation()
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():  # Only add non-empty paragraphs
                    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                    
                    # Add text box with the paragraph content
                    left = top = width = height = Pt(100)
                    width = prs.slide_width - Pt(200)
                    height = prs.slide_height - Pt(200)
                    
                    textbox = slide.shapes.add_textbox(Pt(100), Pt(100), width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = paragraph.text
                    text_frame.word_wrap = True
            
            prs.save(output)
        
        else:
            raise ValueError(f"Unsupported target format: {target_type}")

        output.seek(0)
        return output
    
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"DOCX conversion error: {str(e)}")
