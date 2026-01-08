import io
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import simpleSplit

def convert(file, target_type):
    """
    Convert TXT files to various formats.
    Supports: TXT → DOCX, PDF, PPTX
    """
    try:
        output = io.BytesIO()
        text = file.read().decode('utf-8')

        if target_type.upper() == 'DOCX':
            # Convert TXT to DOCX
            doc = Document()
            
            # Split text into paragraphs and add them
            paragraphs = text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    doc.add_paragraph(para.strip())
            
            doc.save(output)

        elif target_type.upper() == 'PDF':
            # Convert TXT to PDF with proper text wrapping
            c = canvas.Canvas(output, pagesize=letter)
            width, height = letter
            
            y_position = height - 50
            c.setFont("Courier", 10)
            
            for line in text.splitlines():
                # Split long lines to fit page width
                wrapped_lines = simpleSplit(line if line else " ", "Courier", 10, width - 100)
                
                for wrapped_line in wrapped_lines:
                    if y_position < 50:  # Start new page if needed
                        c.showPage()
                        c.setFont("Courier", 10)
                        y_position = height - 50
                    
                    c.drawString(50, y_position, wrapped_line)
                    y_position -= 12
            
            c.save()

        elif target_type.upper() == 'PPTX':
            # Convert TXT to PPTX
            prs = Presentation()
            
            # Split text into chunks (one per slide, max ~500 chars)
            max_chars_per_slide = 500
            paragraphs = text.split('\n\n')
            current_slide_text = ""
            
            for para in paragraphs:
                if len(current_slide_text) + len(para) > max_chars_per_slide and current_slide_text:
                    # Create slide with current text
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    left = top = Pt(50)
                    width = prs.slide_width - Pt(100)
                    height = prs.slide_height - Pt(100)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = current_slide_text.strip()
                    text_frame.word_wrap = True
                    
                    current_slide_text = para + "\n\n"
                else:
                    current_slide_text += para + "\n\n"
            
            # Add remaining text
            if current_slide_text.strip():
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                left = top = Pt(50)
                width = prs.slide_width - Pt(100)
                height = prs.slide_height - Pt(100)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = current_slide_text.strip()
                text_frame.word_wrap = True
            
            prs.save(output)
        
        else:
            raise ValueError(f"Unsupported target format: {target_type}")

        output.seek(0)
        return output
    
    except Exception as e:
        raise Exception(f"TXT conversion error: {str(e)}")
