import io
import os
import tempfile
import subprocess
import zipfile
from pptx import Presentation
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import simpleSplit
from PIL import Image

def convert(file, target_type):
    """
    Convert PPTX files to various formats.
    Supports: PPTX → PDF, Images (PNG, JPG), DOCX, TXT
    """
    try:
        output = io.BytesIO()
        prs = Presentation(file)

        if target_type.upper() == 'PDF':
            # Create PDF with better formatting
            c = canvas.Canvas(output, pagesize=letter)
            width, height = letter
            
            for slide_num, slide in enumerate(prs.slides):
                # Extract text from all shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text)
                
                # Write slide number
                c.setFont("Helvetica-Bold", 14)
                c.drawString(50, height - 50, f"Slide {slide_num + 1}")
                
                # Write text content
                c.setFont("Helvetica", 11)
                y_position = height - 80
                
                for text in slide_text:
                    # Split long lines
                    lines = simpleSplit(text, "Helvetica", 11, width - 100)
                    for line in lines:
                        if y_position < 50:  # Start new page if needed
                            c.showPage()
                            y_position = height - 50
                        c.drawString(50, y_position, line)
                        y_position -= 15
                    y_position -= 10  # Extra space between shapes
                
                c.showPage()
            
            c.save()

        elif target_type.upper() in ['PNG', 'JPG', 'JPEG']:
            # Convert PPTX to images using LibreOffice
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx:
                # Save the presentation to temp file
                file.seek(0)
                temp_pptx.write(file.read())
                temp_pptx_path = temp_pptx.name
            
            try:
                # Create temporary output directory
                temp_dir = tempfile.mkdtemp()
                
                # Convert to PDF first using LibreOffice
                subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', temp_dir, temp_pptx_path
                ], check=True, capture_output=True, timeout=30)
                
                # Find the generated PDF
                pdf_filename = os.path.splitext(os.path.basename(temp_pptx_path))[0] + '.pdf'
                pdf_path = os.path.join(temp_dir, pdf_filename)
                
                # Convert PDF to images
                from pdf2image import convert_from_path
                images = convert_from_path(pdf_path, dpi=200)
                
                if len(images) == 1:
                    # Single slide
                    format_map = {'JPG': 'JPEG', 'JPEG': 'JPEG', 'PNG': 'PNG'}
                    img_format = format_map.get(target_type.upper(), 'PNG')
                    images[0].save(output, format=img_format, quality=95)
                else:
                    # Multiple slides - create ZIP
                    zip_buffer = io.BytesIO()
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                        for idx, image in enumerate(images):
                            img_buffer = io.BytesIO()
                            format_map = {'JPG': 'JPEG', 'JPEG': 'JPEG', 'PNG': 'PNG'}
                            img_format = format_map.get(target_type.upper(), 'PNG')
                            ext = 'jpg' if target_type.upper() in ['JPG', 'JPEG'] else 'png'
                            
                            image.save(img_buffer, format=img_format, quality=95)
                            img_buffer.seek(0)
                            zip_file.writestr(f'slide_{idx + 1}.{ext}', img_buffer.read())
                    
                    zip_buffer.seek(0)
                    output.write(zip_buffer.read())
                
                # Cleanup
                os.remove(pdf_path)
                os.rmdir(temp_dir)
            finally:
                os.remove(temp_pptx_path)

        elif target_type.upper() == 'DOCX':
            # Convert PPTX to DOCX
            doc = Document()
            
            for slide_num, slide in enumerate(prs.slides):
                # Add slide header
                doc.add_heading(f'Slide {slide_num + 1}', level=2)
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        doc.add_paragraph(shape.text)
                
                # Add page break between slides (except last one)
                if slide_num < len(prs.slides) - 1:
                    doc.add_page_break()
            
            doc.save(output)

        elif target_type.upper() == 'TXT':
            # Convert PPTX to TXT
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"=== Slide {slide_num + 1} ===\n\n"
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text += shape.text + "\n\n"
                text += "\n"
            
            output.write(text.encode('utf-8'))
        
        else:
            raise ValueError(f"Unsupported target format: {target_type}")

        output.seek(0)
        return output
    
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"PPTX conversion error: {str(e)}")
