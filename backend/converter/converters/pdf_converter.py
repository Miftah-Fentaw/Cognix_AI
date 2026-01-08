import io
import zipfile
from PyPDF2 import PdfReader
import pdfplumber
from pdf2image import convert_from_bytes
from docx import Document
from pptx import Presentation
from pptx.util import Pt

def convert(file, target_type):
    """
    Convert PDF files to various formats.
    Supports: PDF → DOCX, TXT, Images (PNG, JPG, WEBP), PPTX
    """
    try:
        output = io.BytesIO()
        content = file.read()

        if target_type.upper() == 'TXT':
            # Extract text from PDF
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                text = ""
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
            output.write(text.encode('utf-8'))

        elif target_type.upper() == 'DOCX':
            # Convert PDF to DOCX
            doc = Document()
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        if i > 0:
                            doc.add_page_break()
                        doc.add_paragraph(page_text)
            doc.save(output)

        elif target_type.upper() in ['PNG', 'JPG', 'JPEG', 'WEBP']:
            # Convert PDF pages to images
            images = convert_from_bytes(content, dpi=200)
            
            if len(images) == 1:
                # Single page - return as single image
                format_map = {'JPG': 'JPEG', 'JPEG': 'JPEG', 'PNG': 'PNG', 'WEBP': 'WEBP'}
                img_format = format_map.get(target_type.upper(), 'PNG')
                images[0].save(output, format=img_format, quality=95)
            else:
                # Multiple pages - create a ZIP file with all images
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                    for idx, image in enumerate(images):
                        img_buffer = io.BytesIO()
                        format_map = {'JPG': 'JPEG', 'JPEG': 'JPEG', 'PNG': 'PNG', 'WEBP': 'WEBP'}
                        img_format = format_map.get(target_type.upper(), 'PNG')
                        ext = target_type.lower()
                        
                        image.save(img_buffer, format=img_format, quality=95)
                        img_buffer.seek(0)
                        zip_file.writestr(f'page_{idx + 1}.{ext}', img_buffer.read())
                
                zip_buffer.seek(0)
                output.write(zip_buffer.read())

        elif target_type.upper() == 'PPTX':
            # Convert PDF to PPTX
            prs = Presentation()
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                    page_text = page.extract_text()
                    
                    if page_text:
                        # Add text box
                        left = top = Pt(50)
                        width = prs.slide_width - Pt(100)
                        height = prs.slide_height - Pt(100)
                        
                        textbox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = textbox.text_frame
                        text_frame.text = page_text
                        text_frame.word_wrap = True
            
            prs.save(output)
        
        else:
            raise ValueError(f"Unsupported target format: {target_type}")

        output.seek(0)
        return output
    
    except Exception as e:
        raise Exception(f"PDF conversion error: {str(e)}")
