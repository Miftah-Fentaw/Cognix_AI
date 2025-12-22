import PyPDF2
from docx import Document
from pptx import Presentation

def extract_text_from_file(file):
    """
    Accepts uploaded file (PDF, DOCX, PPTX)
    Returns extracted text as string
    """
    filename = file.name.lower()
    text = ""

    if filename.endswith(".pdf"):
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"

    elif filename.endswith(".docx"):
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif filename.endswith(".pptx"):
        ppt = Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    else:
        raise ValueError("Unsupported file type")

    return text
